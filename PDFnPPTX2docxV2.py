import os
import sys
from pptx import Presentation
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from PIL import Image
import pdfplumber
import fitz  # PyMuPDF

def extract_text_runs(shape):
    content = []
    if not shape.has_text_frame:
        return content

    for para in shape.text_frame.paragraphs:
        para_data = {
            "is_bullet": para.level > 0 or para.text.strip().startswith("•"),
            "runs": []
        }
        for run in para.runs:
            font = run.font
            color = None
            if font.color and hasattr(font.color, "rgb"):
                try:
                    rgb = font.color.rgb
                    if rgb:
                        color = RGBColor(rgb[0], rgb[1], rgb[2])
                except:
                    color = None

            para_data["runs"].append({
                "text": run.text,
                "bold": font.bold,
                "italic": font.italic,
                "underline": font.underline,
                "size": font.size.pt if font.size else None,
                "name": font.name,
                "color": color
            })
        content.append(para_data)
    return content

def add_text_to_docx(doc, text_content):
    for para in text_content:
        doc_para = doc.add_paragraph()
        if para["is_bullet"]:
            doc_para.style = 'List Bullet'
        for run in para["runs"]:
            doc_run = doc_para.add_run(run["text"])
            doc_run.bold = run["bold"]
            doc_run.italic = run["italic"]
            doc_run.underline = run["underline"]
            if run["size"]:
                doc_run.font.size = Pt(run["size"])
            if run["name"]:
                doc_run.font.name = run["name"]
            if run["color"]:
                doc_run.font.color.rgb = run["color"]

def process_shape(doc, shape, slide_num, image_count, image_dir):
    if hasattr(shape, "visible") and not shape.visible:
        return image_count

    if shape.shape_type == 6 and hasattr(shape, "shapes"):
        for subshape in shape.shapes:
            image_count = process_shape(doc, subshape, slide_num, image_count, image_dir)
        return image_count

    if shape.has_text_frame:
        text_content = extract_text_runs(shape)
        add_text_to_docx(doc, text_content)

    elif shape.shape_type == 13 and hasattr(shape, "image"):
        image = shape.image
        image_format = image.ext
        filename = f"{image_dir}/image_{slide_num}_{image_count}.{image_format}"
        with open(filename, "wb") as f:
            f.write(image.blob)
        try:
            img = Image.open(filename)
            width_inches = min(5, img.size[0] / 96)
        except:
            width_inches = 3
        doc.add_picture(filename, width=Inches(width_inches))
        image_count += 1

    elif shape.has_table:
        table = shape.table
        doc_table = doc.add_table(rows=len(table.rows), cols=len(table.columns))
        doc_table.style = "Table Grid"
        for r_idx, row in enumerate(table.rows):
            for c_idx, cell in enumerate(row.cells):
                doc_table.cell(r_idx, c_idx).text = cell.text_frame.text

    elif shape.shape_type == 3 and "Chart" in str(shape.name):
        doc.add_paragraph("[Chart Placeholder: Cannot render chart image]")

    elif "SmartArt" in str(shape.name):
        doc.add_paragraph("[SmartArt Placeholder]")
        if shape.has_text_frame:
            text_content = extract_text_runs(shape)
            add_text_to_docx(doc, text_content)

    return image_count

def pptx_to_docx(pptx_path, docx_path, image_dir="pptx_images"):
    prs = Presentation(pptx_path)
    doc = Document()

    if not os.path.exists(image_dir):
        os.makedirs(image_dir)

    image_count = 0

    for i, slide in enumerate(prs.slides):
        doc.add_heading(f"Slide {i + 1}", level=1)
        for shape in slide.shapes:
            image_count = process_shape(doc, shape, i + 1, image_count, image_dir)
        doc.add_page_break()

    doc.save(docx_path)
    print(f"[✓] Converted PPTX '{pptx_path}' to '{docx_path}' with images saved in '{image_dir}/'.")

def pdf_to_docx(pdf_path, docx_path, image_dir="pdf_images"):
    if not os.path.exists(image_dir):
        os.makedirs(image_dir)

    doc = Document()
    pdf = fitz.open(pdf_path)

    for i, page in enumerate(pdf):
        doc.add_heading(f"Page {i + 1}", level=1)

        # Extract and add text
        text = page.get_text()
        for line in text.split("\n"):
            doc.add_paragraph(line)

        # Extract and add images
        images = page.get_images(full=True)
        for j, img in enumerate(images):
            xref = img[0]
            pix = fitz.Pixmap(pdf, xref)
            img_path = f"{image_dir}/page_{i + 1}_img_{j + 1}.png"
            if pix.n < 5:
                pix.save(img_path)
            else:
                pix1 = fitz.Pixmap(fitz.csRGB, pix)
                pix1.save(img_path)
            doc.add_picture(img_path, width=Inches(4))
        doc.add_page_break()

    doc.save(docx_path)
    print(f"[✓] Converted PDF '{pdf_path}' to '{docx_path}' with images saved in '{image_dir}/'.")

def convert_to_docx(input_path, output_path):
    ext = os.path.splitext(input_path)[-1].lower()
    if ext == ".pptx":
        pptx_to_docx(input_path, output_path)
    elif ext == ".pdf":
        pdf_to_docx(input_path, output_path)
    else:
        print("Unsupported file type. Only .pptx and .pdf are allowed.")
        sys.exit(1)

if __name__ == "__main__":
    input_path = input("Enter the path to the input .pptx or .pdf file: ").strip()
    output_path = input("Enter the desired output .docx filename (leave blank for default): ").strip()

    # If output_path is empty, derive it from input_path
    if not output_path:
        output_path = os.path.splitext(input_path)[0] + ".docx"

    convert_to_docx(input_path, output_path)
